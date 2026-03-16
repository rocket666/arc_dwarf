#!/usr/bin/env python3
"""从 copilot_training_guide.md 生成 PPT 培训幻灯片。"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- 配色方案 ----------
BG_DARK = RGBColor(0x1E, 0x1E, 0x2E)      # 深色背景
BG_SECTION = RGBColor(0x18, 0x18, 0x28)    # 章节页背景
ACCENT = RGBColor(0x58, 0x9C, 0xFF)        # 蓝色强调
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)       # 绿色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x99, 0x99, 0x99)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
CODE_BG = RGBColor(0x28, 0x28, 0x3C)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  space_before=6, space_after=2, font_name="Microsoft YaHei", level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    return p


def add_code_box(slide, left, top, width, height, code_text, font_size=11):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    # 背景填充
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT2
        p.font.name = "Consolas"
        p.space_before = Pt(1)
        p.space_after = Pt(1)
    return tf


def add_table_slide(slide, left, top, width, rows_data, col_widths=None, font_size=12):
    """rows_data: list of lists; first row is header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.4 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Microsoft YaHei"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = LIGHT_GRAY
            # Header row bg
            if r == 0:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0x30, 0x30, 0x50)
            else:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0x24, 0x24, 0x38) if r % 2 == 1 else RGBColor(0x2A, 0x2A, 0x42)
    return table


def make_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_SECTION)
    # 左边装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_text_box(slide, 0.8, 2.5, 8.5, 1.2, title, font_size=36, color=ACCENT, bold=True)
    if subtitle:
        add_text_box(slide, 0.8, 3.7, 8.5, 0.8, subtitle, font_size=18, color=DIM_GRAY)
    return slide


def make_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x24)
    shape.line.fill.background()
    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.88), Inches(10), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    add_text_box(slide, 0.5, 0.15, 9, 0.65, title, font_size=24, color=WHITE, bold=True)
    return slide


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ===== 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 1.5, 1.5, 10, 1.5,
                 "使用 GitHub Copilot 开发\narc_dwarf 工具链",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 3.5, 10, 0.7,
                 "培训指南 — 从零到完整项目的全流程演示",
                 font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 5.0, 10, 0.5,
                 "嵌入式汽车雷达系统 · 自动化标定工具链 · ELF/DWARF + A2L",
                 font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 6.2, 10, 0.5,
                 "2026年3月",
                 font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # ===== 目录 =====
    slide = make_content_slide(prs, "目录")
    items = [
        "1. 背景介绍 — arc_dwarf 是什么、解决什么问题",
        "2. 环境准备 — VS Code + Copilot + Python 配置",
        "3. 完整开发流程 — 6 个 Phase 从零到可用",
        "4. 碰到的问题与解决方案",
        "5. 注意事项 — 沟通技巧、安全、规范",
        "6. 后续使用建议 — 日常场景、配置优化、进阶用法",
    ]
    tf = add_text_box(slide, 1.0, 1.3, 11, 5.5, "", font_size=20)
    tf.paragraphs[0].text = items[0]
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].space_after = Pt(8)
    for item in items[1:]:
        add_paragraph(tf, item, font_size=20, color=LIGHT_GRAY, space_before=12, space_after=8)

    # ===== 第1章：背景介绍 =====
    make_section_slide(prs, "1. 背景介绍", "arc_dwarf 是什么、Copilot 承担什么角色")

    # --- 什么是 arc_dwarf ---
    slide = make_content_slide(prs, "1.1 什么是 arc_dwarf")
    tf = add_text_box(slide, 0.8, 1.2, 11.5, 5, "", font_size=18)
    tf.paragraphs[0].text = "面向嵌入式汽车雷达系统的自动化标定工具链"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ ELF/DWARF 解析 — 提取结构体定义、成员偏移、类型信息",
        "▸ 符号表查询 — 获取全局变量的运行时物理地址",
        "▸ A2L 标定文件解析 — 解析 ASAP2 标定变量定义",
        "▸ 成员地址计算 — 沿嵌套路径自动计算深层成员的物理地址",
    ]:
        add_paragraph(tf, item, font_size=18, color=LIGHT_GRAY, space_before=10)

    add_paragraph(tf, "", space_before=10)
    add_paragraph(tf, "典型查询路径：", font_size=16, color=DIM_GRAY, space_before=10)
    add_code_box(slide, 1.0, 4.8, 11, 0.7,
                 "gRspCfg.RspFastTimeCmbCtrlParam.FftCtrl.cfg_fft_win_addr", font_size=14)

    # --- Copilot 角色 ---
    slide = make_content_slide(prs, "1.2 Copilot 在本项目中的角色")
    add_table_slide(slide, 0.8, 1.3, 11.5, [
        ["角色", "具体工作"],
        ["代码开发者", "实现 resolve_member_offset()、get_member_address() 等核心 API"],
        ["调试助手", "定位文件路径、解析 MAP 文件、运行验证脚本"],
        ["文档撰写者", "更新 API 文档和快速上手指南"],
        ["运维助手", "执行 git commit/push 等版本管理操作"],
        ["问题解决者", "发现并修复 ModuleNotFoundError 等运行时问题"],
    ], col_widths=[2.5, 9], font_size=14)

    # ===== 第2章：环境准备 =====
    make_section_slide(prs, "2. 环境准备", "VS Code + Copilot + Python 开发环境配置")

    slide = make_content_slide(prs, "2.1 必要工具与配置")
    add_table_slide(slide, 0.8, 1.3, 5.5, [
        ["工具", "说明"],
        ["VS Code", "主要开发环境"],
        ["GitHub Copilot 扩展", "Copilot + Copilot Chat"],
        ["Python 3.9+", "运行环境"],
        ["Git", "版本管理"],
    ], col_widths=[2.5, 3], font_size=14)
    tf = add_text_box(slide, 6.8, 1.3, 6, 4.5, "", font_size=16)
    tf.paragraphs[0].text = "推荐配置"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ 启用 Copilot Chat Agent 模式",
        "▸ 确保 pyproject.toml 在工作区根目录",
        "▸ 打开相关文件（A2L/源码）作为上下文",
    ]:
        add_paragraph(tf, item, font_size=15, color=LIGHT_GRAY, space_before=10)

    add_code_box(slide, 6.8, 4.0, 6, 2.2, """# 项目初始化
git clone https://github.com/rocket666/arc_dwarf.git
cd arc_dwarf
pip install -e ".[dev]"
python -c "from arc_dwarf import DwarfTypeExplorer; print('OK')"
""", font_size=12)

    # --- 数据文件 ---
    slide = make_content_slide(prs, "2.2 准备数据文件")
    tf = add_text_box(slide, 0.8, 1.3, 11.5, 5, "", font_size=18)
    tf.paragraphs[0].text = "将以下文件放在项目根目录："
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ CPD_BPM_App_Core1.elf — 带 DWARF 调试信息的 ELF（编译时加 -g）",
        "▸ CPD_BPM_App_Core1.map — 链接器生成的 MAP 文件",
        "▸ Langcang_bpm_solution.a2l — A2L 标定文件",
    ]:
        add_paragraph(tf, item, font_size=16, color=LIGHT_GRAY, space_before=10)
    add_paragraph(tf, "", space_before=16)
    add_paragraph(tf, "⚠  .elf 和 .map 已在 .gitignore 中配置为不跟踪，需手动拷贝", font_size=15, color=ORANGE, space_before=10)

    # ===== 第3章：完整开发流程 =====
    make_section_slide(prs, "3. 完整开发流程", "从零开始 — 6 个 Phase 的全程演示")

    # --- Phase 1 ---
    slide = make_content_slide(prs, "Phase 1：创建 GitHub 仓库")
    tf = add_text_box(slide, 0.8, 1.3, 11.5, 2, "", font_size=18)
    tf.paragraphs[0].text = "Commit: 75d5149  Initial commit"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = DIM_GRAY
    tf.paragraphs[0].font.name = "Consolas"
    add_paragraph(tf, "", space_before=10)
    add_paragraph(tf, "在 GitHub 网页端创建空仓库 rocket666/arc_dwarf", font_size=20, color=LIGHT_GRAY, space_before=8)
    add_paragraph(tf, "自动生成 README.md — 仅包含项目名", font_size=20, color=LIGHT_GRAY, space_before=8)
    add_paragraph(tf, "", space_before=16)
    add_paragraph(tf, "💡  这是唯一不经过 Copilot 的步骤", font_size=18, color=ACCENT, bold=True, space_before=10)
    add_paragraph(tf, "后续所有开发均在 VS Code + Copilot Agent 模式下完成", font_size=16, color=DIM_GRAY, space_before=6)

    # --- Phase 2 ---
    slide = make_content_slide(prs, "Phase 2：构建完整项目工程（1/3）")
    tf = add_text_box(slide, 0.8, 1.2, 11.5, 1.5, "", font_size=16)
    tf.paragraphs[0].text = "Commit: 03ddc37  refactor: migrate to modern Python project structure"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = DIM_GRAY
    tf.paragraphs[0].font.name = "Consolas"
    add_paragraph(tf, "", space_before=6)
    add_paragraph(tf, "用户指令：", font_size=16, color=ACCENT, bold=True, space_before=4)
    add_code_box(slide, 0.8, 2.5, 11.5, 0.6,
                 "将现有代码迁移到现代 Python 项目结构（src-layout），添加完整的项目配置、测试、文档和 CI", font_size=13)
    add_paragraph(tf, "", space_before=6)
    tf2 = add_text_box(slide, 0.8, 3.5, 11.5, 3.5, "", font_size=18)
    tf2.paragraphs[0].text = "Copilot 一次性完成 — 29 个文件、1822 行代码："
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ 项目配置：pyproject.toml (PEP 621 + hatchling)、.gitignore、LICENSE (MIT)",
        "▸ 核心源码：6 个模块 (dwarf_explorer / a2l_parser / cli / type_dumper / model / textout)",
        "▸ 测试套件：14 个单元测试、pytest fixtures、sample.a2l 测试数据",
        "▸ 文档骨架：Sphinx 配置 + 3 篇 Markdown 文档",
        "▸ CI/CD：GitHub Actions 4 版本矩阵测试 + PyPI 自动发布",
    ]:
        add_paragraph(tf2, item, font_size=15, color=LIGHT_GRAY, space_before=8)

    # Phase 2 (2/3) - 项目结构
    slide = make_content_slide(prs, "Phase 2：项目结构一览（2/3）")
    add_code_box(slide, 0.5, 1.2, 5.8, 5.8, """arc_dwarf/
├── .github/workflows/
│   ├── ci.yml          # CI 4版本矩阵
│   └── release.yml     # PyPI 自动发布
├── .gitignore
├── LICENSE (MIT)
├── pyproject.toml      # PEP 621 配置
├── docs/
│   ├── getting_started.md
│   ├── api_reference.md
│   └── architecture.md
├── examples/
│   ├── basic_usage.py
│   └── interactive_tool.py
├── src/arc_dwarf/
│   ├── dwarf_explorer.py  # 412行 核心
│   ├── a2l_parser.py
│   ├── cli.py
│   ├── type_dumper.py
│   ├── model.py
│   └── textout.py
└── tests/
    ├── test_a2l_parser.py
    ├── test_dwarf_explorer.py
    ├── test_model.py
    └── test_type_dumper.py""", font_size=12)
    add_table_slide(slide, 6.8, 1.2, 6, [
        ["设计决策", "选择", "原因"],
        ["构建后端", "hatchling", "轻量, PEP 621 原生"],
        ["项目布局", "src-layout", "社区最佳实践"],
        ["Lint", "ruff", "速度快, 替代 flake8"],
        ["类型检查", "mypy", "Python 标准"],
        ["许可证", "MIT", "宽松开源"],
    ], col_widths=[1.5, 1.8, 2.7], font_size=13)

    # Phase 2 (3/3) - 模块关系
    slide = make_content_slide(prs, "Phase 2：核心模块关系（3/3）")
    add_code_box(slide, 1.0, 1.3, 5, 2.5, """A2lParser ──────────┐
                    ├──▶ CLI (交互式工具)
DwarfTypeExplorer ──┤
       │            │
       ▼            │
  TypeTreeDumper ───┘
       │
       ▼
  TextEmitter (输出格式化)""", font_size=14)
    tf = add_text_box(slide, 6.5, 1.3, 6.3, 5.5, "", font_size=16)
    tf.paragraphs[0].text = "dwarf_explorer.py（412 行）— 核心引擎"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ _index_types() — 建立类型名索引",
        "▸ resolve_typedefs_and_qualifiers() — typedef 链",
        "▸ type_byte_size() — 递归计算类型大小",
        "▸ member_offset() — 解码 DW_AT_data_member_location",
        "▸ parse_struct_or_union_members() — 解析成员",
        "▸ get_variable_info() — 符号表查找",
    ]:
        add_paragraph(tf, item, font_size=14, color=LIGHT_GRAY, space_before=8)

    add_paragraph(tf, "", space_before=14)
    add_paragraph(tf, "💡 Copilot 自动选择了社区最佳实践", font_size=16, color=ACCENT, bold=True, space_before=6)
    add_paragraph(tf, "用户只给了一条指令，全部结构和配置由 Copilot 决定", font_size=14, color=DIM_GRAY, space_before=4)

    # --- Phase 3 ---
    slide = make_content_slide(prs, "Phase 3：需求驱动开发 — 成员地址解析 API（1/3）")
    tf = add_text_box(slide, 0.8, 1.2, 11.5, 1, "", font_size=14)
    tf.paragraphs[0].text = "Commit: 377a764  feat: add member path resolution API"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = DIM_GRAY
    tf.paragraphs[0].font.name = "Consolas"
    add_text_box(slide, 0.8, 1.9, 6, 0.5, "用户指令：", font_size=16, color=ACCENT, bold=True)
    add_code_box(slide, 0.8, 2.5, 11.5, 0.9,
                 "使用 CPD_BPM_App_Core1.elf 和 CPD_BPM_App_Core1.map，根据数据结构 Rsp_Params_t，\n计算 gRspCfg.RspFastTimeCmbCtrlParam.FftCtrl.cfg_fft_size 的物理地址", font_size=13)
    tf2 = add_text_box(slide, 0.8, 3.8, 11.5, 3.2, "", font_size=16)
    tf2.paragraphs[0].text = "Copilot 处理步骤："
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for i, item in enumerate([
        "① 环境感知 — 搜索文件系统定位 .elf/.map，子代理阅读全部源码",
        "② 从 MAP 提取基地址 — 搜索 gRspCfg → 发现 0x20FD1630",
        "③ 发现能力缺口 — 现有代码无法沿路径计算偏移",
        "④ 自主实现新 API — 修改 6 个文件，新增 304 行代码",
        "⑤ 安装并验证 — pip install -e . → 运行脚本确认结果",
    ]):
        add_paragraph(tf2, item, font_size=15, color=LIGHT_GRAY, space_before=8)

    # Phase 3 (2/3) - 实现细节
    slide = make_content_slide(prs, "Phase 3：文件变更与输出结果（2/3）")
    add_table_slide(slide, 0.5, 1.2, 6.5, [
        ["修改文件", "变更内容"],
        ["dwarf_explorer.py", "+141行 resolve_member_offset() / get_member_address()"],
        ["model.py", "+12行 MemberAddress dataclass"],
        ["__init__.py", "+2行 导出 MemberAddress"],
        ["getting_started.md", "+48行 使用示例"],
        ["api_reference.md", "+55行 API 文档"],
        ["test_dwarf_explorer.py", "+46行 测试用例"],
    ], col_widths=[2.5, 4], font_size=13)
    add_code_box(slide, 7.5, 1.2, 5.3, 3.5, """运行结果：

变量名:       gRspCfg
成员路径:     ...FftCtrl.cfg_fft_size
基地址:       0x20FD16B0
成员偏移:     0x0054 (84 bytes)
物理地址:     0x20FD1704
成员类型:     uint16
成员大小:     2 bytes""", font_size=13)
    tf = add_text_box(slide, 7.5, 5.0, 5.3, 1.8, "", font_size=14)
    tf.paragraphs[0].text = "⚠ 关键发现"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf, "A2L 地址: 0x20F80000 (占位)", font_size=14, color=LIGHT_GRAY, space_before=6)
    add_paragraph(tf, "ELF 实际: 0x20FD16B0 (真实)", font_size=14, color=ACCENT2, space_before=4)
    add_paragraph(tf, "→ 始终以 ELF 符号表为准！", font_size=14, color=ORANGE, space_before=4)

    # Phase 3 (3/3) - 后续操作
    slide = make_content_slide(prs, "Phase 3：后续操作链（3/3）")
    tf = add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", font_size=16)
    tf.paragraphs[0].text = "连续对话 — 无缝衔接多个任务"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        '用户："cfg_fft_win_addr 的地址计算，并给出命令行"',
        '   → Copilot 直接调用 API，输出 0x20FD1700，附命令行',
        '',
        '用户："在 docs 中增加获取成员路径的 API 使用说明"',
        '   → Copilot 同时更新 getting_started.md 和 api_reference.md',
        '',
        '用户："自动 commit 仓库中的更改"',
        '   → Copilot: git status → git add -A → git commit（自动生成 message）',
        '',
        '用户："push 到 web 端"',
        '   → Copilot 先征求确认，确认后执行 git push origin main',
    ]:
        color = ACCENT if item.startswith('用户') else (DIM_GRAY if item.startswith('   →') else LIGHT_GRAY)
        fs = 15 if item.startswith('用户') else 14
        bold = item.startswith('用户')
        add_paragraph(tf, item, font_size=fs, color=color, bold=bold, space_before=4)

    # --- Phase 4 & 5 ---
    slide = make_content_slide(prs, "Phase 4 & 5：示例创建与问题修复")
    # Phase 4
    tf = add_text_box(slide, 0.5, 1.2, 6, 2.8, "", font_size=16)
    tf.paragraphs[0].text = "Phase 4：创建可运行示例"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf, 'cce805c  examples: add member_address.py', font_size=12, color=DIM_GRAY, space_before=6)
    add_paragraph(tf, '', space_before=4)
    add_paragraph(tf, '▸ 用户："examples 下增加可运行的 case"', font_size=14, color=LIGHT_GRAY, space_before=4)
    add_paragraph(tf, '▸ Copilot 创建脚本 → 自动运行验证 → commit + push', font_size=14, color=LIGHT_GRAY, space_before=4)
    # Phase 5
    tf2 = add_text_box(slide, 0.5, 4.2, 6, 3, "", font_size=16)
    tf2.paragraphs[0].text = "Phase 5：修复跨平台问题"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf2, 'cadc554  fix: add sys.path fallback', font_size=12, color=DIM_GRAY, space_before=6)
    add_paragraph(tf2, '', space_before=4)
    add_paragraph(tf2, '▸ Windows 报错: ModuleNotFoundError', font_size=14, color=ORANGE, space_before=4)
    add_paragraph(tf2, '▸ 用户粘贴 Traceback → Copilot 诊断并修复', font_size=14, color=LIGHT_GRAY, space_before=4)
    add_paragraph(tf2, '▸ 添加 sys.path 回退兼容未安装场景', font_size=14, color=LIGHT_GRAY, space_before=4)
    # 右侧流程图
    add_code_box(slide, 7.0, 1.2, 5.8, 5.5, """完整开发流程回顾：

Phase 1  GitHub 创建空仓库
   ↓
Phase 2  Copilot 生成完整项目
         (29 文件, 1822 行)
   ↓
Phase 3  需求驱动: 实现地址计算 API
         (6 文件, 304 行)
   ↓
Phase 4  创建可运行示例
   ↓
Phase 5  修复跨平台问题
   ↓
Phase 6  编写培训文档

━━━━━━━━━━━━━━━━━━━━━━━
总计: 5 commits, 2200+ 行
用户输入: ~10 条自然语言
用户手写代码: 0 行""", font_size=13)

    # ===== 第4章：碰到的问题 =====
    make_section_slide(prs, "4. 碰到的问题与解决", "开发过程中遇到的 4 个典型问题")

    slide = make_content_slide(prs, "问题 1 & 2")
    # 问题1
    tf = add_text_box(slide, 0.5, 1.2, 6, 2.8, "", font_size=16)
    tf.paragraphs[0].text = "问题 1：A2L 地址与 ELF 实际地址不一致"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf, "A2L: 0x20F80000 (占位)  vs  ELF: 0x20FD16B0 (真实)", font_size=14, color=LIGHT_GRAY, space_before=10)
    add_paragraph(tf, "", space_before=6)
    add_paragraph(tf, "✅ 解决：始终用 get_variable_info() 从 ELF 符号表获取", font_size=14, color=ACCENT2, space_before=4)
    add_paragraph(tf, "📌 教训：A2L 地址仅供参考，以 ELF 为准", font_size=14, color=ACCENT, space_before=4)
    # 问题2
    tf2 = add_text_box(slide, 0.5, 4.2, 6, 3, "", font_size=16)
    tf2.paragraphs[0].text = "问题 2：未安装包时 import 失败"
    tf2.paragraphs[0].font.size = Pt(17)
    tf2.paragraphs[0].font.color.rgb = ORANGE
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf2, "ModuleNotFoundError: No module named 'arc_dwarf'", font_size=13, color=LIGHT_GRAY, space_before=10)
    add_paragraph(tf2, "", space_before=6)
    add_paragraph(tf2, "✅ 解决：添加 sys.path 回退 + pip install -e .", font_size=14, color=ACCENT2, space_before=4)
    add_paragraph(tf2, "📌 教训：示例需兼顾已安装和未安装两种场景", font_size=14, color=ACCENT, space_before=4)

    # 问题3&4
    tf3 = add_text_box(slide, 7.0, 1.2, 5.8, 2.8, "", font_size=16)
    tf3.paragraphs[0].text = "问题 3：ELF/MAP 未纳入版本控制"
    tf3.paragraphs[0].font.size = Pt(17)
    tf3.paragraphs[0].font.color.rgb = ORANGE
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf3, ".gitignore 排除了 *.elf *.map（大文件）", font_size=14, color=LIGHT_GRAY, space_before=10)
    add_paragraph(tf3, "", space_before=6)
    add_paragraph(tf3, "✅ 解决：文档说明需手动拷贝", font_size=14, color=ACCENT2, space_before=4)
    add_paragraph(tf3, "📌 建议：考虑 Git LFS 或团队共享存储", font_size=14, color=ACCENT, space_before=4)

    tf4 = add_text_box(slide, 7.0, 4.2, 5.8, 3, "", font_size=16)
    tf4.paragraphs[0].text = "问题 4：Copilot 操作确认机制"
    tf4.paragraphs[0].font.size = Pt(17)
    tf4.paragraphs[0].font.color.rgb = ORANGE
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf4, "git push 前 Copilot 主动请求确认", font_size=14, color=LIGHT_GRAY, space_before=10)
    add_paragraph(tf4, "", space_before=6)
    add_paragraph(tf4, "✅ 这是安全特性，不可逆操作需人工确认", font_size=14, color=ACCENT2, space_before=4)
    add_paragraph(tf4, "📌 涉及远程操作、文件删除等保持确认", font_size=14, color=ACCENT, space_before=4)

    # ===== 第5章：注意事项 =====
    make_section_slide(prs, "5. 注意事项", "沟通技巧、安全规范、项目约定")

    slide = make_content_slide(prs, "5.1 与 Copilot 沟通的技巧")
    add_table_slide(slide, 0.5, 1.3, 12, [
        ["技巧", "说明", "示例"],
        ["打开相关文件", "Copilot 自动读取编辑器中的文件", "打开 A2L 后提问，无需解释变量"],
        ["选中关键文本", "选中代码段可以聚焦问题", "选中 gRspCfg 后提问地址计算"],
        ["分步提问", "复杂任务分解为小步骤", "先计算 → 再写文档 → 再 commit"],
        ["明确动作", "避免模糊指令", "「提交 commit 并 push」而非「保存一下」"],
        ["粘贴错误", "完整输出比口头描述更高效", "粘贴完整 Traceback"],
    ], col_widths=[2, 4.5, 5.5], font_size=13)

    slide = make_content_slide(prs, "5.2 安全、规范与 ELF 要求")
    tf = add_text_box(slide, 0.5, 1.2, 6, 3, "", font_size=16)
    tf.paragraphs[0].text = "安全与权限"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ 不可逆操作需确认（push / reset / delete）",
        "▸ 不要粘贴密钥、密码等敏感信息",
        "▸ Copilot 生成的代码仍需人工审查",
        "▸ 确保大文件和敏感文件不被提交",
    ]:
        add_paragraph(tf, item, font_size=15, color=LIGHT_GRAY, space_before=8)

    tf2 = add_text_box(slide, 0.5, 4.5, 6, 2.5, "", font_size=16)
    tf2.paragraphs[0].text = "项目规范"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ Conventional Commits: feat: / fix: / docs:",
        "▸ src-layout: 源码 src/ 测试 tests/ 示例 examples/",
        "▸ type hints + dataclass + context manager",
    ]:
        add_paragraph(tf2, item, font_size=15, color=LIGHT_GRAY, space_before=8)

    tf3 = add_text_box(slide, 7.0, 1.2, 5.8, 3, "", font_size=16)
    tf3.paragraphs[0].text = "ELF 文件要求"
    tf3.paragraphs[0].font.size = Pt(18)
    tf3.paragraphs[0].font.color.rgb = ACCENT
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "▸ 编译时必须加 -g 保留 DWARF 信息",
        "▸ 不要 strip 去除符号表",
        "▸ ELF 与 MAP 需来自同一次编译",
    ]:
        add_paragraph(tf3, item, font_size=15, color=LIGHT_GRAY, space_before=8)

    # ===== 第6章：后续建议 =====
    make_section_slide(prs, "6. 后续使用建议", "日常场景、模式选择、配置优化、进阶用法")

    slide = make_content_slide(prs, "6.1 日常开发场景")
    add_table_slide(slide, 0.5, 1.3, 12, [
        ["场景", "Copilot 用法示例"],
        ["查询变量地址", "「计算 gSolFlowSts.xxx 的物理地址」"],
        ["添加新类型支持", "「在 A2L 解析器中增加对 AXIS_PTS 块的支持」"],
        ["编写单元测试", "「为 resolve_member_offset 添加数组下标测试」"],
        ["代码重构", "「将 interactive_shell 的查询逻辑抽取为独立函数」"],
        ["批量计算", "「遍历所有 MEASUREMENT 变量，输出真实物理地址」"],
    ], col_widths=[3, 9], font_size=14)

    slide = make_content_slide(prs, "6.2 Agent 模式 vs Chat 模式 vs Inline 补全")
    add_table_slide(slide, 0.5, 1.3, 12, [
        ["模式", "适用场景", "特点"],
        ["Agent 模式", "多步骤任务（写码+测试+文档+commit）", "执行命令、读写文件、管理 Git"],
        ["Chat 模式", "快速问答、代码片段生成", "只提供建议，不直接操作"],
        ["Inline 补全", "编写代码时的自动补全", "根据上下文实时建议"],
    ], col_widths=[2.5, 4.5, 5], font_size=14)
    tf = add_text_box(slide, 0.5, 4.5, 12, 1.5, "", font_size=18)
    tf.paragraphs[0].text = "💡 推荐：arc_dwarf 日常使用优先选择 Agent 模式"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf, "因为它可以直接运行 Python 脚本、执行终端命令并返回结果", font_size=15, color=LIGHT_GRAY, space_before=6)

    slide = make_content_slide(prs, "6.3 提升效果的配置 & 6.4 进阶用法")
    tf = add_text_box(slide, 0.5, 1.2, 6, 3, "", font_size=16)
    tf.paragraphs[0].text = "Copilot 指令文件"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    add_paragraph(tf, "创建 .github/copilot-instructions.md：", font_size=14, color=LIGHT_GRAY, space_before=8)
    add_code_box(slide, 0.5, 2.5, 6, 2.5, """# Copilot Instructions
- 本项目是嵌入式雷达标定工具链
- ELF/A2L 在项目根目录
- 地址以 ELF 符号表为准
- Conventional Commits 格式
- PEP 8 规范""", font_size=12)

    tf2 = add_text_box(slide, 7.0, 1.2, 5.8, 5.5, "", font_size=16)
    tf2.paragraphs[0].text = "进阶用法建议"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for item in [
        "① 批量自动化",
        "   批量计算所有 A2L 变量地址，自动回填",
        "② CI 集成",
        "   GitHub Actions 地址一致性检查",
        "③ CLI 非交互模式",
        "   命令行参数直接指定变量和路径",
        "④ 自定义 Copilot Skill",
        "   创建 SKILL.md 实现一句话触发查询",
        "⑤ Memory 功能",
        "   Copilot 记住常用变量名和构建命令",
    ]:
        bold = item[0] in "①②③④⑤"
        color = LIGHT_GRAY if not bold else ACCENT2
        add_paragraph(tf2, item, font_size=14, color=color, bold=bold, space_before=4 if not bold else 10)

    # ===== 结束页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 1.5, 2.0, 10, 1.2,
                 "谢谢！",
                 font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 3.5, 10, 1.2,
                 "GitHub Copilot + arc_dwarf\n让标定工具开发更高效",
                 font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_table_slide(slide, 3.0, 5.0, 7, [
        ["指标", "数据"],
        ["Copilot 参与的 commit", "4 / 5"],
        ["总代码行数", "2,200+"],
        ["用户手写代码", "0 行"],
    ], col_widths=[3, 4], font_size=14)

    # ---------- 保存 ----------
    output_path = "/workspaces/arc_dwarf/docs/copilot_training_guide.pptx"
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    build_ppt()
